"""
Python library-based converters (fallback methods)
"""

from datetime import datetime
from pathlib import Path

from .base_converter import BaseConverter
from ..utils import setup_logger
from config.settings import (
    RAG_OPTIMIZATION_ENABLED,
    EXCEL_TABLE_MAX_ROWS_PER_PAGE,
    EXCEL_TABLE_OPTIMIZATION,
    CITATION_INCLUDE_FILENAME,
    CITATION_INCLUDE_PAGE,
    CITATION_INCLUDE_DATE,
    CITATION_DATE_FORMAT,
)


class PythonLibraryConverter(BaseConverter):
    """Base class for Python library converters"""

    def __init__(self):
        self.logger = setup_logger(__name__)
        super().__init__()
    
    def is_available(self) -> bool:
        """Check if required libraries are available"""
        return True

    def _apply_pdf_metadata(self, pdf_path: Path, source: Path, page_count: int | None = None) -> None:
        if not RAG_OPTIMIZATION_ENABLED:
            return
        try:
            from pypdf import PdfReader, PdfWriter
        except ImportError:
            self.logger.debug("pypdf not installed; skipping metadata for %s", source.name)
            return

        try:
            reader = PdfReader(str(pdf_path))
            actual_pages = page_count or len(reader.pages)
            writer = PdfWriter()
            for page in reader.pages:
                writer.add_page(page)

            metadata_parts = []
            keywords = [source.stem, 'PythonConverter', 'RAG']

            if CITATION_INCLUDE_FILENAME:
                metadata_parts.append(f"Source: {source.name}")
                keywords.append(source.name)
            if CITATION_INCLUDE_PAGE and actual_pages:
                metadata_parts.append(f"Pages: {actual_pages}")
                keywords.append(f"pages:{actual_pages}")
            if CITATION_INCLUDE_DATE:
                metadata_parts.append(
                    f"Converted: {datetime.now().strftime(CITATION_DATE_FORMAT)}"
                )

            metadata = {
                '/Title': source.stem,
                '/Author': 'AI4Team Converter',
                '/Subject': 'RAG Export',
                '/Creator': 'Python Library Converter',
                '/Producer': 'ReportLab' if pdf_path.suffix.lower() == '.pdf' else 'Python',
                '/Keywords': ', '.join(dict.fromkeys(keywords)),
            }
            if metadata_parts:
                metadata['/Comments'] = ' | '.join(metadata_parts)

            existing_metadata = dict(reader.metadata or {})
            existing_metadata.update(metadata)
            writer.add_metadata(existing_metadata)

            temp_path = pdf_path.parent / (pdf_path.name + '.tmp')
            with open(temp_path, 'wb') as temp_file:
                writer.write(temp_file)
            temp_path.replace(pdf_path)
        except Exception as metadata_error:
            self.logger.debug("Failed to inject metadata for %s: %s", source.name, metadata_error)


class DocxConverter(PythonLibraryConverter):
    """Convert DOCX using docx2pdf"""
    
    def convert(self, input_file: Path, output_file: Path) -> bool:
        try:
            from docx2pdf import convert
            convert(str(input_file), str(output_file))
            if output_file.exists():
                self._apply_pdf_metadata(output_file, input_file)
            return True
        except Exception:
            return False


class XlsxConverter(PythonLibraryConverter):
    """Convert XLSX using openpyxl and reportlab"""
    
    def convert(self, input_file: Path, output_file: Path) -> bool:
        try:
            from openpyxl import load_workbook
            from reportlab.lib.pagesizes import letter, landscape
            from reportlab.lib import colors
            from reportlab.platypus import (
                SimpleDocTemplate,
                Table,
                TableStyle,
                PageBreak,
                Paragraph,
            )
            from reportlab.lib.styles import getSampleStyleSheet
            
            # Load Excel workbook
            wb = load_workbook(input_file, data_only=True)
            
            # Create PDF
            doc = SimpleDocTemplate(
                str(output_file),
                pagesize=landscape(letter),
                title=input_file.stem,
            )
            elements = []
            styles = getSampleStyleSheet()
            
            # Process each sheet
            sheet_names = wb.sheetnames
            for sheet_index, sheet_name in enumerate(sheet_names):
                ws = wb[sheet_name]
                
                # Extract data
                data = []
                for row in ws.iter_rows(values_only=True):
                    row_data = [str(cell) if cell is not None else '' for cell in row]
                    data.append(row_data)
                
                if data:
                    if RAG_OPTIMIZATION_ENABLED:
                        citation_bits = []
                        if CITATION_INCLUDE_FILENAME:
                            citation_bits.append(input_file.name)
                        citation_bits.append(f"Sheet: {sheet_name}")
                        if CITATION_INCLUDE_DATE:
                            citation_bits.append(datetime.now().strftime(CITATION_DATE_FORMAT))
                        header_text = " | ".join(citation_bits)
                        elements.append(Paragraph(f"<b>{header_text}</b>", styles['Normal']))

                    header = data[0]
                    rows = data[1:] if len(data) > 1 else []
                    use_table_chunks = (
                        RAG_OPTIMIZATION_ENABLED
                        and EXCEL_TABLE_OPTIMIZATION
                        and EXCEL_TABLE_MAX_ROWS_PER_PAGE > 0
                    )
                    if use_table_chunks and rows:
                        chunks = [
                            rows[i:i + EXCEL_TABLE_MAX_ROWS_PER_PAGE]
                            for i in range(0, len(rows), EXCEL_TABLE_MAX_ROWS_PER_PAGE)
                        ]
                    else:
                        chunks = [rows]

                    for chunk_index, chunk in enumerate(chunks):
                        table_data = [header] + chunk if header else chunk
                        table = Table(table_data)
                        if header:
                            try:
                                table.repeatRows = 1
                            except Exception:
                                pass

                        style = TableStyle([
                            ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                            ('FONTSIZE', (0, 0), (-1, 0), 10),
                            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                            ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                            ('GRID', (0, 0), (-1, -1), 0.5, colors.black)
                        ])
                        table.setStyle(style)
                        elements.append(table)
                        if chunk_index < len(chunks) - 1:
                            elements.append(PageBreak())

                if sheet_index < len(sheet_names) - 1:
                    elements.append(PageBreak())
            
            # Build PDF
            doc.build(elements)
            if output_file.exists():
                self._apply_pdf_metadata(output_file, input_file)
            return True
            
        except Exception:
            return False


class PptxConverter(PythonLibraryConverter):
    """Convert PPTX using python-pptx and reportlab"""
    
    def convert(self, input_file: Path, output_file: Path) -> bool:
        try:
            from pptx import Presentation
            from reportlab.lib.pagesizes import letter, landscape
            from reportlab.pdfgen import canvas
            from reportlab.lib.units import inch
            
            # Load presentation
            prs = Presentation(input_file)
            
            # Create PDF
            c = canvas.Canvas(str(output_file), pagesize=landscape(letter))
            width, height = landscape(letter)
            
            # Process each slide
            for slide_num, slide in enumerate(prs.slides, 1):
                # Extract text from slide
                y_position = height - inch
                
                c.setFont("Helvetica-Bold", 16)
                c.drawString(inch, y_position, f"Slide {slide_num}")
                y_position -= 0.5 * inch
                
                c.setFont("Helvetica", 12)
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text = shape.text.strip()
                        
                        # Simple text wrapping
                        for line in text.split('\n'):
                            if y_position < inch:
                                c.showPage()
                                c.setFont("Helvetica", 12)
                                y_position = height - inch
                            
                            # Truncate if too long
                            if len(line) > 100:
                                line = line[:97] + "..."
                            
                            c.drawString(inch, y_position, line)
                            y_position -= 0.3 * inch
                
                if RAG_OPTIMIZATION_ENABLED and CITATION_INCLUDE_FILENAME:
                    footer = input_file.name
                    if CITATION_INCLUDE_DATE:
                        footer += f" | {datetime.now().strftime(CITATION_DATE_FORMAT)}"
                    c.setFont("Helvetica", 8)
                    if hasattr(c, "drawRightString"):
                        c.drawRightString(width - inch, 0.5 * inch, footer)
                    else:
                        c.drawString(width - inch, 0.5 * inch, footer)
                c.showPage()
            
            c.save()
            if output_file.exists():
                self._apply_pdf_metadata(output_file, input_file)
            return True
            
        except Exception:
            return False


class CsvConverter(PythonLibraryConverter):
    """Convert CSV to PDF using reportlab"""
    
    def convert(self, input_file: Path, output_file: Path) -> bool:
        try:
            import csv
            from reportlab.lib.pagesizes import letter, landscape
            from reportlab.lib import colors
            from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, PageBreak
            from reportlab.lib.styles import getSampleStyleSheet
            
            # Read CSV file
            with open(input_file, 'r', encoding='utf-8', errors='ignore') as f:
                reader = csv.reader(f)
                data = list(reader)
            
            if not data:
                return False
            
            # Create PDF
            doc = SimpleDocTemplate(
                str(output_file),
                pagesize=landscape(letter),
                title=input_file.stem,
            )
            elements = []
            
            # Add title
            styles = getSampleStyleSheet()
            title = Paragraph(f"<b>{input_file.name}</b>", styles['Title'])
            elements.append(title)
            
            # Create table
            if RAG_OPTIMIZATION_ENABLED:
                citation_bits = []
                if CITATION_INCLUDE_FILENAME:
                    citation_bits.append(input_file.name)
                if CITATION_INCLUDE_DATE:
                    citation_bits.append(datetime.now().strftime(CITATION_DATE_FORMAT))
                header_text = " | ".join(citation_bits)
                elements.append(Paragraph(f"<b>{header_text}</b>", styles['Normal']))

            header = data[0]
            rows = data[1:] if len(data) > 1 else []
            max_rows = (
                EXCEL_TABLE_MAX_ROWS_PER_PAGE
                if (RAG_OPTIMIZATION_ENABLED and EXCEL_TABLE_OPTIMIZATION)
                else 0
            )
            if max_rows and rows:
                row_chunks = [
                    rows[i:i + max_rows]
                    for i in range(0, len(rows), max_rows)
                ]
            else:
                row_chunks = [rows]

            for idx, chunk in enumerate(row_chunks):
                table_data = [header] + chunk if header else chunk
                table = Table(table_data)
                if header:
                    try:
                        table.repeatRows = 1
                    except Exception:
                        pass
                
                style = TableStyle([
                    ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                    ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                    ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                    ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                    ('FONTSIZE', (0, 0), (-1, 0), 10),
                    ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                    ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                    ('GRID', (0, 0), (-1, -1), 1, colors.black),
                    ('FONTSIZE', (0, 0), (-1, -1), 8),
                ])
                table.setStyle(style)
                
                elements.append(table)
                if idx < len(row_chunks) - 1:
                    elements.append(PageBreak())
            
            # Build PDF
            doc.build(elements)
            if output_file.exists():
                self._apply_pdf_metadata(output_file, input_file)
            return True
            
        except Exception:
            return False



